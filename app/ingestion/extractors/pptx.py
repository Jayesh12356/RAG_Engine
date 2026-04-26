"""PowerPoint (.pptx) extractor.

Each slide becomes one ``ParsedPage`` containing the title, every shape's text
(bullets/runs preserved), embedded table-shape contents rendered as a GFM
table, and finally the speaker notes appended under a "Notes" sub-heading.
"""
from __future__ import annotations

import os

import structlog

from app.ingestion.pdf_parser import ParsedPage

logger = structlog.get_logger(__name__)


def _render_pptx_table(table) -> str:
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
    if not rows:
        return ""
    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    out: list[str] = []
    out.append("| " + " | ".join(header) + " |")
    out.append("| " + " | ".join(["---"] * len(header)) + " |")
    for r in body:
        cells = list(r) + [""] * max(0, len(header) - len(r))
        out.append("| " + " | ".join(cells[: len(header)]) + " |")
    return "\n".join(out)


def parse_pptx(path: str) -> list[ParsedPage]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("python-pptx is required for .pptx ingestion") from exc

    pdf_name = os.path.basename(path)
    service_name = os.path.splitext(pdf_name)[0].replace("_", " ")
    pres = Presentation(path)

    pages: list[ParsedPage] = []
    total = max(1, len(pres.slides))
    has_any_table = False
    for idx, slide in enumerate(pres.slides, start=1):
        title = ""
        body_parts: list[str] = []

        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""

        for shape in slide.shapes:
            if shape == getattr(slide.shapes, "title", None):
                continue
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                tf_text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text and p.text.strip()
                ).strip()
                if tf_text:
                    body_parts.append(tf_text)
            if getattr(shape, "has_table", False):
                rendered = _render_pptx_table(shape.table)
                if rendered:
                    body_parts.append(rendered)
                    has_any_table = True

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""
        if notes:
            body_parts.append(f"## Notes\n{notes}")

        title_line = f"# {title}\n\n" if title else ""
        text = (title_line + "\n\n".join(body_parts)).strip()
        if not text:
            continue

        pages.append(
            ParsedPage(
                page_number=idx,
                text=text,
                pdf_name=pdf_name,
                service_name=service_name,
                section_title=title or f"Slide {idx}",
                total_pages=total,
                kind="slide",
            )
        )

    if not pages:
        pages.append(
            ParsedPage(
                page_number=1,
                text="(empty presentation)",
                pdf_name=pdf_name,
                service_name=service_name,
                section_title="Presentation",
                total_pages=1,
                kind="prose",
            )
        )

    logger.info("pptx.parse.complete", path=path, slides=len(pages), tables=has_any_table)
    return pages
